"""
send_weekly_shell_consumos.py — Reporte semanal Shell al área comercial.
========================================================================

Cada lunes envía a jperez, pgomez, dhevia, jcaceres, wsoto:
  - Excel con TODAS las OTs Shell del MES en curso (hasta el domingo previo)
  - PPT con métricas de la SEMANA ISO ANTERIOR (lun→dom)

Hilo Gmail: por MES CALENDARIO del lunes de envío.
  Lunes 27-jul → nuevo hilo "Julio 2026"
  Lunes 03-ago → nuevo hilo "Agosto 2026" (aunque el reporte cubra sem 31)
  Lunes 10, 17, 24, 31-ago → reply al hilo de agosto
  Lunes 07-sep → nuevo hilo "Septiembre 2026"

Persistencia del thread en Supabase (tabla email_threads).

Ejecución:
    python send_weekly_shell_consumos.py
    python send_weekly_shell_consumos.py --dry-run
    python send_weekly_shell_consumos.py --test-email jgavidia@occimiano.cl
    python send_weekly_shell_consumos.py --fecha 2026-08-03   (simular lunes)
    python send_weekly_shell_consumos.py --force-new-thread
"""

from __future__ import annotations

import argparse
import io
import os
import smtplib
import sys
import traceback
from datetime import datetime, date, timedelta, timezone
from email.message import EmailMessage
from email.utils import make_msgid, formatdate
from pathlib import Path
from typing import Optional

import pandas as pd
import requests


# ── Cargar .env ─────────────────────────────────────────────────────────────
def _load_env():
    env_path = Path(__file__).parent / ".env"
    if env_path.exists():
        for line in env_path.read_text(encoding="utf-8").splitlines():
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                k, v = line.split("=", 1)
                os.environ.setdefault(k.strip(), v.strip())


_load_env()

GMAIL_SENDER   = "jgavidia@occimiano.cl"
GMAIL_DISPLAY  = "Operaciones Occimiano"
GMAIL_PASSWORD = os.getenv("GMAIL_APP_PASSWORD", "")

SUPABASE_URL = os.getenv("SUPABASE_URL", "")
SUPABASE_KEY = os.getenv("SUPABASE_KEY", "")

DESTINATARIOS = [
    "jperez@occimiano.cl",
    "pgomez@occimiano.cl",
    "dhevia@occimiano.cl",
    "jcaceres@occimiano.cl",
    "wsoto@occimiano.cl",
]

TOPIC = "shell_consumos"

MESES_ES = ["", "Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
            "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"]


def log(msg: str):
    print(f"[{datetime.now().strftime('%H:%M:%S')}] {msg}", flush=True)


# ── Supabase helpers ────────────────────────────────────────────────────────
def _sb_get(path: str, params: dict, limit_pag: int = 1000) -> list:
    h = {"apikey": SUPABASE_KEY, "Authorization": f"Bearer {SUPABASE_KEY}"}
    out = []
    offset = 0
    while True:
        p = dict(params)
        p["limit"] = limit_pag
        p["offset"] = offset
        r = requests.get(f"{SUPABASE_URL}/rest/v1/{path}", headers=h,
                         params=p, timeout=30)
        r.raise_for_status()
        batch = r.json() or []
        out.extend(batch)
        if len(batch) < limit_pag:
            break
        offset += limit_pag
    return out


def _sb_upsert(path: str, row: dict, on_conflict: str) -> bool:
    h = {"apikey": SUPABASE_KEY, "Authorization": f"Bearer {SUPABASE_KEY}",
         "Content-Type": "application/json",
         "Prefer": "resolution=merge-duplicates,return=minimal"}
    r = requests.post(f"{SUPABASE_URL}/rest/v1/{path}",
                      params={"on_conflict": on_conflict},
                      headers=h, json=[row], timeout=30)
    return r.status_code in (200, 201, 204)


def _get_thread(topic: str, mes: str) -> Optional[dict]:
    try:
        rows = _sb_get("email_threads", {
            "select": "topic,mes,message_id,subject",
            "topic":  f"eq.{topic}",
            "mes":    f"eq.{mes}",
        })
        return rows[0] if rows else None
    except requests.exceptions.HTTPError as e:
        if e.response is not None and e.response.status_code in (404, 403, 400):
            log(f"  ⚠ Tabla email_threads no accesible ({e.response.status_code}) — tratando como nuevo hilo.")
            return None
        raise


def _save_thread(topic: str, mes: str, message_id: str, subject: str) -> bool:
    try:
        return _sb_upsert("email_threads", {
            "topic": topic, "mes": mes, "message_id": message_id,
            "subject": subject,
            "sent_at": datetime.now(timezone.utc).isoformat(),
        }, on_conflict="topic,mes")
    except Exception as e:
        log(f"  ⚠ No se pudo guardar thread ({e}) — el próximo envío abrirá nuevo hilo.")
        return False


# ── Data loading ────────────────────────────────────────────────────────────
def load_shell_data(mes_yyyy_mm: str, semana_ini: date, semana_fin: date,
                    corte_hasta: date) -> tuple:
    """Devuelve (df_mes, df_sem):
       - df_mes: OTs Shell del mes con fecha <= corte_hasta (ayer inclusive)
       - df_sem: subset con creation_date dentro de la semana ISO anterior
    """
    year_str, month_str = mes_yyyy_mm.split("-")
    inicio_mes = f"{year_str}-{month_str}-01"

    # OTs preventivas Shell del mes (tabla ordenes_trabajo)
    ots = _sb_get("ordenes_trabajo", {
        "select": ("id_ot,codigo_activo,nombre_activo,ubicacion,"
                   "fecha_creacion,responsable,cliente,estacion,codigo_eds,"
                   "tipo_tarea"),
        "cliente":        "eq.SHELL (Enex)",
        "fecha_creacion": f"gte.{inicio_mes}",
    })
    df_ot = pd.DataFrame(ots)
    if df_ot.empty:
        return df_ot, df_ot

    # Filtrar preventivas (tipo_tarea contiene PREVENTIVA)
    df_ot = df_ot[df_ot["tipo_tarea"].fillna("").str.upper().str.contains("PREVENTIVA", na=False)].copy()
    if df_ot.empty:
        return df_ot, df_ot

    # Fecha limpia
    df_ot["_fecha_dt"] = pd.to_datetime(df_ot["fecha_creacion"], errors="coerce", utc=True)\
                          .dt.tz_convert("America/Santiago")
    df_ot["_fecha"] = df_ot["_fecha_dt"].dt.date

    folios = list(df_ot["id_ot"].dropna().unique())
    if not folios:
        return df_ot.iloc[0:0], df_ot.iloc[0:0]

    # numerales_subtarea por lote (in.(...))
    subs = []
    for i in range(0, len(folios), 80):
        chunk = folios[i:i+80]
        subs.extend(_sb_get("numerales_subtarea", {
            "select": ("id_ot,codigo_activo,nombre_activo,tipo_activo,"
                       "numeral_inicial,numeral_final,fichas_periodo,"
                       "bomba_dosificadora,lts_hr_produccion_final,"
                       "consumo_cera_pct,consumo_shampoo_pct,consumo_cepillo_pct,"
                       "tiempo_fichas_seg,cubre_fichero,form_tiene_bomba,"
                       "fecha_inicio_subtarea,fecha_fin_subtarea"),
            "id_ot": f"in.({','.join(chunk)})",
        }))
    df_sub = pd.DataFrame(subs)

    # Estación (nombre + comuna)
    est = _sb_get("estaciones_servicio", {
        "select": "eds_occim,nombre,comuna,cliente",
        "cliente": "eq.SHELL (Enex)",
    })
    df_est = pd.DataFrame(est)
    if not df_est.empty:
        est_map = df_est.set_index(df_est["eds_occim"].astype(str))
    else:
        est_map = pd.DataFrame()

    # Un registro por OT: preferimos lavadora con form_tiene_bomba=True
    rows_out = []
    for _, ot in df_ot.iterrows():
        id_ot = ot["id_ot"]
        sub_ot = df_sub[df_sub["id_ot"] == id_ot] if not df_sub.empty else pd.DataFrame()
        lav = sub_ot[sub_ot["tipo_activo"] == "lavadora"] if not sub_ot.empty else pd.DataFrame()
        if len(lav) > 1 and "form_tiene_bomba" in lav.columns:
            lav = lav.sort_values("form_tiene_bomba", ascending=False)
        asp = sub_ot[sub_ot["tipo_activo"] == "aspiradora"] if not sub_ot.empty else pd.DataFrame()

        def _rp(df, col, default=None):
            if df.empty or col not in df.columns: return default
            v = df.iloc[0][col]
            if v is None or (isinstance(v, float) and pd.isna(v)): return default
            s = str(v).strip()
            return default if s in ("", "None", "nan", "null") else s

        def _rp_hora(df, col):
            if df.empty or col not in df.columns: return None
            vals = pd.to_datetime(df[col], errors="coerce").dropna()
            if vals.empty: return None
            ts = vals.min() if "inicio" in col else vals.max()
            try:
                if ts.tz is not None: ts = ts.tz_convert("America/Santiago")
                return ts.time()
            except Exception:
                return None

        eds_code = str(ot.get("codigo_eds", "") or "")
        est_row = est_map.loc[eds_code] if (not est_map.empty and eds_code in est_map.index) else None
        nombre_est = est_row["nombre"] if est_row is not None else (ot.get("estacion") or "—")
        ciudad = est_row["comuna"].title() if (est_row is not None and pd.notna(est_row.get("comuna"))) else "—"

        rows_out.append({
            "Fecha":             ot["_fecha_dt"],
            "N° OT":             id_ot,
            "Equipo":            _rp(lav, "codigo_activo") or _rp(sub_ot, "codigo_activo") or "—",
            "Código EDS":        eds_code or "—",
            "Nombre Estación":   nombre_est,
            "Ciudad":            ciudad,
            "Hora inicio":       _rp_hora(sub_ot, "fecha_inicio_subtarea"),
            "Hora término":      _rp_hora(sub_ot, "fecha_fin_subtarea"),
            "Técnico":           ot.get("responsable") or "—",
            "Anterior":          _rp(lav, "numeral_inicial", "—"),
            "Actual":            _rp(lav, "numeral_final", "—"),
            "Anterior 2":        _rp(asp, "numeral_inicial", "—"),
            "Actual 2":          _rp(asp, "numeral_final", "—"),
            "Lavado":            _rp(lav, "fichas_periodo", "—"),
            "Aspirado":          _rp(asp, "fichas_periodo", "—"),
            "Tipo de bomba":     _rp(lav, "bomba_dosificadora"),
            "Producción Lts/hr": _rp(lav, "lts_hr_produccion_final"),
            "Cera %":            _fmt_pct(_rp(lav, "consumo_cera_pct")),
            "Shampoo %":         _fmt_pct(_rp(lav, "consumo_shampoo_pct")),
            "Cepillo %":         _fmt_pct(_rp(lav, "consumo_cepillo_pct")),
            "Lavado (seg)":      _rp(lav, "tiempo_fichas_seg"),
            "Cubre Fichero":     _fmt_cf(_rp(lav, "cubre_fichero")),
            "_fecha_dt":         ot["_fecha_dt"],
        })

    df_mes = pd.DataFrame(rows_out).sort_values("_fecha_dt", ascending=False).reset_index(drop=True)

    # Excluir el día actual: solo hasta 'corte_hasta' inclusive (ayer)
    df_mes = df_mes[df_mes["_fecha_dt"].dt.date <= corte_hasta].copy().reset_index(drop=True)

    # Subset semana ISO anterior
    df_sem = df_mes[df_mes["_fecha_dt"].dt.date.between(semana_ini, semana_fin)].copy()
    return df_mes, df_sem


def _fmt_pct(v):
    if v is None: return None
    try:
        n = float(str(v).rstrip("%").replace(",", "."))
        if n <= 1: n *= 100
        return round(n / 100.0, 4)   # queda como 0.30 → Excel lo formatea 30%
    except (ValueError, TypeError):
        return None


def _fmt_cf(v):
    if v is None: return None
    s = str(v).strip().upper()
    if s in ("SI", "SÍ", "TRUE", "YES"): return "✅ Sí"
    if s in ("NO", "FALSE"): return "❌ No"
    return v


# ── Excel builder ───────────────────────────────────────────────────────────
def build_excel(df_mes: pd.DataFrame) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill
    wb = Workbook()
    ws = wb.active
    ws.title = "Datos_bomba"

    cols = ["Fecha","N° OT","Equipo","Código EDS","Nombre Estación","Ciudad",
            "Hora inicio","Hora término","Técnico","Anterior","Actual",
            "Anterior 2","Actual 2","Lavado","Aspirado","Tipo de bomba",
            "Producción Lts/hr","Cera %","Shampoo %","Cepillo %",
            "Lavado (seg)","Cubre Fichero"]

    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill("solid", fgColor="1F4E78")
    center = Alignment(horizontal="center", vertical="center")

    for c_i, col in enumerate(cols, 1):
        cell = ws.cell(1, c_i, col)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center

    for r_i, (_, row) in enumerate(df_mes.iterrows(), 2):
        for c_i, col in enumerate(cols, 1):
            val = row.get(col)
            if col == "Fecha" and pd.notna(val):
                cell = ws.cell(r_i, c_i, val.to_pydatetime().replace(tzinfo=None))
                cell.number_format = "dd/mm/yyyy"
            elif col in ("Hora inicio", "Hora término"):
                cell = ws.cell(r_i, c_i, val)
                cell.number_format = "hh:mm"
            elif col in ("Cera %", "Shampoo %", "Cepillo %") and val is not None:
                cell = ws.cell(r_i, c_i, val)
                cell.number_format = "0.0%"
            else:
                ws.cell(r_i, c_i, val)

    # Auto width aprox
    widths = {"A":11,"B":10,"C":9,"D":10,"E":32,"F":14,"G":10,"H":10,"I":25,
              "J":10,"K":10,"L":10,"M":10,"N":10,"O":10,"P":22,"Q":16,"R":9,
              "S":9,"T":9,"U":12,"V":12}
    for col, w in widths.items():
        ws.column_dimensions[col].width = w
    ws.freeze_panes = "A2"

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── PPT builder (usa la plantilla del área comercial) ──────────────────────
_PPT_TEMPLATE = Path(__file__).parent / "assets" / "plantilla_shell_consumos.pptx"


def _replace_text(shape, new_text: str, color=None) -> None:
    """Reemplaza el texto conservando el formato del primer run.
    Si `new_text` contiene '\\n' crea múltiples párrafos.
    Si `color` es RGBColor, se aplica al texto (útil para alertas rojas)."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return

    lines = new_text.split("\n") if new_text else [""]
    p0 = tf.paragraphs[0]
    # Vaciar todos los runs del primer párrafo excepto el primero
    if p0.runs:
        p0.runs[0].text = lines[0]
        if color is not None:
            p0.runs[0].font.color.rgb = color
        for r in p0.runs[1:]:
            r.text = ""
    else:
        run = p0.add_run()
        run.text = lines[0]
        if color is not None:
            run.font.color.rgb = color

    # Vaciar párrafos existentes adicionales
    for p2 in tf.paragraphs[1:]:
        for r in p2.runs:
            r.text = ""

    # Agregar párrafos nuevos para líneas adicionales
    for line in lines[1:]:
        p_new = tf.add_paragraph()
        # Copiar formato del primer párrafo (best-effort)
        run = p_new.add_run()
        run.text = line
        try:
            if p0.runs and p0.runs[0].font.size:
                run.font.size = p0.runs[0].font.size
            if color is not None:
                run.font.color.rgb = color
            elif p0.runs and p0.runs[0].font.color and p0.runs[0].font.color.type is not None:
                run.font.color.rgb = p0.runs[0].font.color.rgb
        except Exception:
            pass


def _walk_named(shapes):
    """Itera recursivamente shapes (entra en GROUPs) devolviendo (name, shape)."""
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            yield from _walk_named(sh.shapes)
        else:
            yield sh.name, sh


def build_ppt(df_sem: pd.DataFrame, mes_lbl: str, sem_iso: int,
              semana_ini: date, semana_fin: date,
              corte_hasta: date | None = None) -> tuple[bytes, list[str]]:
    """Devuelve (pptx_bytes, resumen_lines) para reutilizar el resumen
    también en el cuerpo del correo.

    Nota: `df_sem` puede ser en realidad `df_mes` (consolidado del mes)
    — el nombre se mantiene por compatibilidad interna.
    """
    from pptx import Presentation

    if not _PPT_TEMPLATE.exists():
        raise RuntimeError(f"No se encuentra plantilla PPT: {_PPT_TEMPLATE}")

    total = len(df_sem)

    # ── Cálculos comunes ──
    def _pob(col):
        if col not in df_sem.columns: return 0
        return int(df_sem[col].apply(lambda v: v is not None
                                     and not (isinstance(v,float) and pd.isna(v))
                                     and str(v).strip() not in ("","—","None","nan")).sum())

    def _fmt_pct1(n, d):
        return f"{100*n/d:.1f} %" if d else "—"

    def _fmt_pct1_short(n, d):
        return f"{100*n/d:.1f}%" if d else "—"

    n_bomba   = _pob("Tipo de bomba")
    n_prod    = _pob("Producción Lts/hr")
    n_cera    = _pob("Cera %")
    n_shp     = _pob("Shampoo %")
    n_cep     = _pob("Cepillo %")
    n_lavseg  = _pob("Lavado (seg)")

    n_cf_si = int(df_sem["Cubre Fichero"].fillna("").astype(str).str.contains("Sí", na=False).sum()) if "Cubre Fichero" in df_sem.columns else 0
    n_cf_no = int(df_sem["Cubre Fichero"].fillna("").astype(str).str.contains("❌", na=False).sum()) if "Cubre Fichero" in df_sem.columns else 0
    n_cf_pob = n_cf_si + n_cf_no
    n_cf_sd = total - n_cf_pob

    def _prom_pct(col):
        vals = [v for v in df_sem.get(col, []) if isinstance(v,(int,float)) and not pd.isna(v)]
        if not vals: return None, 0
        return sum(vals)*100/len(vals), len(vals)
    prom_cera_pct, prom_cera_n = _prom_pct("Cera %")
    prom_shp_pct,  prom_shp_n  = _prom_pct("Shampoo %")
    prom_cep_pct,  prom_cep_n  = _prom_pct("Cepillo %")

    def _mas_frec(col):
        s = df_sem.get(col, pd.Series([], dtype=object)).dropna()
        if s.empty: return None, 0
        # Para %s formatear como "30%"
        if col in ("Cera %", "Shampoo %", "Cepillo %"):
            s2 = s.apply(lambda x: f"{int(round(float(x)*100))}%" if isinstance(x,(int,float)) else str(x))
        else:
            s2 = s.astype(str)
        vc = s2.value_counts()
        if vc.empty: return None, 0
        return vc.index[0], int(vc.iloc[0])

    frec_cera_v, frec_cera_n = _mas_frec("Cera %")
    frec_shp_v,  frec_shp_n  = _mas_frec("Shampoo %")
    frec_cep_v,  frec_cep_n  = _mas_frec("Cepillo %")

    # Distribución cera para Text 63
    dist_cera_str = "—"
    if "Cera %" in df_sem.columns:
        s = df_sem["Cera %"].dropna()
        if not s.empty:
            s2 = s.apply(lambda x: f"{int(round(float(x)*100))}%")
            vc = s2.value_counts().head(5)
            dist_cera_str = "Distribución Cera: " + " · ".join(f"{v} ({n})" for v, n in vc.items())

    # Lavado típico (moda) + excepciones (≠ típico) + sin dato
    lav_typ = "—"
    lav_typ_n = 0
    excepciones_lav = []   # (EDS, valor_seg)
    sin_dato_lav    = []   # EDS sin dato
    if "Lavado (seg)" in df_sem.columns and "Código EDS" in df_sem.columns:
        for _, row in df_sem.iterrows():
            v = row["Lavado (seg)"]
            eds = str(row.get("Código EDS") or "").strip() or "—"
            if v is None or (isinstance(v, float) and pd.isna(v)) or str(v).strip() in ("", "—", "None", "nan"):
                sin_dato_lav.append(eds)
                continue
            try:
                n = int(float(v))
                # Los guardamos junto al EDS, después de saber el típico decidimos
                excepciones_lav.append((eds, n))
            except (ValueError, TypeError):
                sin_dato_lav.append(eds)
        # Calcular la moda
        if excepciones_lav:
            vals = [n for _, n in excepciones_lav]
            typ_val = max(set(vals), key=vals.count)
            lav_typ = str(typ_val)
            lav_typ_n = vals.count(typ_val)
            # Solo dejar los que difieren de la moda
            excepciones_lav = [(e, n) for e, n in excepciones_lav if n != typ_val]

    # ── Distribución Tipo de bomba (top 5 + Sin datos) ──
    def _top5_con_sin_datos(col, n_pob, total):
        s = df_sem.get(col, pd.Series([], dtype=object)).dropna().astype(str)
        s = s[s.str.strip() != ""]
        vc = s.value_counts().head(5)
        items = [(str(v)[:40], int(n), 100*n/total if total else 0)
                 for v, n in vc.items()]
        sd = total - n_pob
        if sd > 0:
            items.append(("Sin datos", sd, 100*sd/total if total else 0))
        return items

    tb_items = _top5_con_sin_datos("Tipo de bomba", n_bomba, total)
    while len(tb_items) < 6:
        tb_items.append(("", 0, 0))

    # Producción: top 6 + Sin datos
    prod_items = _top5_con_sin_datos("Producción Lts/hr", n_prod, total)
    while len(prod_items) < 6:
        prod_items.append(("", 0, 0))

    # ── Mapeo textos plantilla ──
    # Rango del mes (desde 01 del mes hasta el corte de datos = ayer)
    _ini_mes = corte_hasta.replace(day=1) if corte_hasta else semana_ini.replace(day=1)
    _fin_rango = corte_hasta if corte_hasta else semana_fin
    _rango_mes = f"{_ini_mes.strftime('%d/%m')} – {_fin_rango.strftime('%d/%m')}"
    subtitulo_s1 = (f"SHELL  ·  {mes_lbl}  ·  {_rango_mes}"
                    f"  ·  {total} registros analizados (consolidado del mes)")
    subtitulo_s2 = (f"SHELL  ·  {mes_lbl}  ·  Distribución de los datos disponibles")

    replacements_s1 = {
        # Header
        "Text 2":  subtitulo_s1,
        # TIPO DE BOMBA
        "Text 6":  f"{n_bomba}/{total}",
        "Text 7":  _fmt_pct1(n_bomba, total),
        "Text 10": f"Faltan {total-n_bomba} registros" if total-n_bomba>0 else "Completo",
        # PRODUCCIÓN LTS/HR
        "Text 13": f"{n_prod}/{total}",
        "Text 14": _fmt_pct1(n_prod, total),
        "Text 17": f"Faltan {total-n_prod} registros" if total-n_prod>0 else "Completo",
        # CONSUMO CERA
        "Text 20": f"{n_cera}/{total}",
        "Text 21": _fmt_pct1(n_cera, total),
        "Text 24": f"Faltan {total-n_cera} registros" if total-n_cera>0 else "Completo",
        # CONSUMO SHAMPOO
        "Text 27": f"{n_shp}/{total}",
        "Text 28": _fmt_pct1(n_shp, total),
        "Text 31": f"Faltan {total-n_shp} registros" if total-n_shp>0 else "Completo",
        # CONSUMO CEPILLO
        "Text 34": f"{n_cep}/{total}",
        "Text 35": _fmt_pct1(n_cep, total),
        "Text 38": f"Faltan {total-n_cep} registros" if total-n_cep>0 else "Completo",
        # LAVADO (SEG)
        "Text 41": f"{n_lavseg}/{total}",
        "Text 42": _fmt_pct1(n_lavseg, total),
        "Text 45": f"Faltan {total-n_lavseg} registros" if total-n_lavseg>0 else "Completo",
        # CUBRE FICHERO
        "Text 48": f"{n_cf_pob}/{total}",
        "Text 49": _fmt_pct1(n_cf_pob, total),
        "Text 52": f"{n_cf_si} Sí · {n_cf_no} No · {n_cf_sd} sin dato",
        # PROMEDIOS OPERATIVOS
        "Text 55": f"{prom_cera_pct:.1f} %" if prom_cera_pct is not None else "—",
        "Text 57": f"{prom_cera_n} registros",
        "Text 58": f"{prom_shp_pct:.1f} %"  if prom_shp_pct  is not None else "—",
        "Text 60": f"{prom_shp_n} registros",
        "Text 61": f"{prom_cep_pct:.1f} %"  if prom_cep_pct  is not None else "—",
        "Text 63": f"{prom_cep_n} registros",
        "Text 64": f"{lav_typ} seg" if lav_typ != "—" else "—",
        "Text 66": f"{lav_typ_n} usan {lav_typ} s" if lav_typ_n else "—",
    }

    # Slide 2: 6 filas de "Tipo de bomba" y 6 filas de "Producción"
    # Cada fila tiene: círculo bullet + nombre + N + %
    tb_labels_texts = [("Shape 5",  "Text 6",  "Text 7",  "Text 8"),
                       ("Shape 9",  "Text 10", "Text 11", "Text 12"),
                       ("Shape 13", "Text 14", "Text 15", "Text 16"),
                       ("Shape 17", "Text 18", "Text 19", "Text 20"),
                       ("Shape 21", "Text 22", "Text 23", "Text 24"),
                       ("Shape 25", "Text 26", "Text 27", "Text 28")]
    prod_labels_texts = [("Text 31","Text 32","Text 33"),
                         ("Text 34","Text 35","Text 36"),
                         ("Text 37","Text 38","Text 39"),
                         ("Text 40","Text 41","Text 42"),
                         ("Text 43","Text 44","Text 45"),
                         ("Text 46","Text 47","Text 48")]

    replacements_s2 = {
        "Text 2":  subtitulo_s2,
        "Text 4":  f"TIPO DE BOMBA  ·  {n_bomba} de {total} registros",
        "Text 30": f"PRODUCCIÓN LTS/HR  ·  {n_prod} de {total} registros",
        # Consumo Insumos
        "Text 51": f"{prom_cera_pct:.1f}%" if prom_cera_pct is not None else "—",
        "Text 53": f"{prom_cera_n} reg.",
        "Text 54": f"Más frecuente: {frec_cera_v} ({frec_cera_n})" if frec_cera_v else "—",
        "Text 55": f"{prom_shp_pct:.1f}%" if prom_shp_pct is not None else "—",
        "Text 57": f"{prom_shp_n} reg.",
        "Text 58": f"Más frecuente: {frec_shp_v} ({frec_shp_n})" if frec_shp_v else "—",
        "Text 59": f"{prom_cep_pct:.1f}%" if prom_cep_pct is not None else "—",
        "Text 61": f"{prom_cep_n} reg.",
        "Text 62": f"Más frecuente: {frec_cep_v} ({frec_cep_n})" if frec_cep_v else "—",
        "Text 63": dist_cera_str,
        # Lavado
        "Text 66": lav_typ,
        "Text 67": f"segundos  ·  {lav_typ_n} reg." if lav_typ_n else "—",
        # Excepciones = OTs con LAVADO ≠ moda (típicamente 210 s)
        "Text 68": f"Excepciones (Sem. {sem_iso}):" if (excepciones_lav or sin_dato_lav) else "",
        "Text 69": ("\n".join(f"{e} → {n} s" for e, n in excepciones_lav[:4])
                    + (f"\n{' · '.join(sin_dato_lav[:3])}{' + ' + str(len(sin_dato_lav)-3) if len(sin_dato_lav)>3 else ''} sin dato"
                       if sin_dato_lav else "")) if (excepciones_lav or sin_dato_lav) else "",
        "Text 70": "",
        "Text 71": "",
        # Cubre fichero
        "Text 74": str(n_cf_si),
        "Text 76": str(n_cf_no),
        "Text 78": str(n_cf_sd),
        "Text 80": _fmt_pct1_short(n_cf_pob, total) + " con info",
    }
    # Guardamos también qué círculos (Shape N) tenemos que OCULTAR porque
    # esa fila quedó vacía (no hay tantos tipos distintos de bomba).
    ocultar_shapes_s2: set[str] = set()
    for i, (shp_bullet, t_n, t_v, t_p) in enumerate(tb_labels_texts):
        nombre, n, pct = tb_items[i]
        replacements_s2[t_n] = nombre or ""
        replacements_s2[t_v] = str(n) if nombre else ""
        replacements_s2[t_p] = f"{pct:.1f}%" if nombre else ""
        if not nombre:
            ocultar_shapes_s2.add(shp_bullet)

    # Filas producción (Text 31..48) — también registramos los vacíos para
    # ocultar el card completo achicándolo, y usar el espacio para el resumen.
    prod_positions = [("Text 31","Text 32","Text 33"),
                      ("Text 34","Text 35","Text 36"),
                      ("Text 37","Text 38","Text 39"),
                      ("Text 40","Text 41","Text 42"),
                      ("Text 43","Text 44","Text 45"),
                      ("Text 46","Text 47","Text 48")]
    filas_prod_activas = 0
    for i, (t_n, t_v, t_p) in enumerate(prod_positions):
        nombre, n, pct = prod_items[i]
        replacements_s2[t_n] = nombre or ""
        replacements_s2[t_v] = str(n) if nombre else ""
        replacements_s2[t_p] = f"{pct:.1f}%" if nombre else ""
        if nombre:
            filas_prod_activas = i + 1

    # ── Resumen automático de alertas para pie de Slide 2 ──────────────
    # (EDS sin cubre fichero, lavado ≠ 210, faltantes por campo, etc.)
    def _eds_sin(col):
        s = df_sem.dropna(subset=["Código EDS"])
        s = s[s[col].isna() | s[col].astype(str).isin(["","—","None","nan"])]
        return s["Código EDS"].astype(str).tolist()

    faltantes_alertas: list[str] = []
    faltantes_por_campo = {
        "Tipo de bomba":      total - n_bomba,
        "Producción Lts/hr":  total - n_prod,
        "Cera %":             total - n_cera,
        "Shampoo %":          total - n_shp,
        "Cepillo %":          total - n_cep,
        "Lavado (seg)":       total - n_lavseg,
        "Cubre fichero":      total - n_cf_pob,
    }
    for campo, nf in faltantes_por_campo.items():
        if nf > 0:
            eds_ls = _eds_sin(campo if campo != "Cubre fichero" else "Cubre Fichero")
            eds_str = f" ({', '.join(eds_ls[:5])}{'…' if len(eds_ls)>5 else ''})" if eds_ls else ""
            faltantes_alertas.append(f"  · {campo}: {nf} sin dato{eds_str}")

    lav_atipico_str = ""
    if excepciones_lav:
        lav_atipico_str = ("  · Lavado atípico (≠ típico " + lav_typ + " s): "
                           + ", ".join(f"{e} → {n} s" for e, n in excepciones_lav[:5]))

    cubre_no_str = ""
    if "Cubre Fichero" in df_sem.columns:
        eds_no = df_sem[df_sem["Cubre Fichero"].fillna("").astype(str).str.contains("❌", na=False)]["Código EDS"].astype(str).tolist()
        if eds_no:
            cubre_no_str = "  · Sin cubre fichero instalado: " + ", ".join(eds_no[:8])

    bombas_raras_str = ""
    if "Tipo de bomba" in df_sem.columns:
        s = df_sem["Tipo de bomba"].dropna().astype(str)
        s = s[s.str.strip() != ""]
        vc = s.value_counts()
        raras = [f"{v} ({n})" for v, n in vc.items() if n == 1]
        if raras:
            bombas_raras_str = "  · Bombas de baja frecuencia (1 vez): " + ", ".join(raras[:5])

    resumen_lines = [
        f"Resumen — {mes_lbl} ({_rango_mes})   ·   {total} registros analizados (mes completo)",
    ]
    if faltantes_alertas:
        resumen_lines.append("Faltantes de datos:")
        resumen_lines.extend(faltantes_alertas[:6])
    if lav_atipico_str:
        resumen_lines.append(lav_atipico_str)
    if cubre_no_str:
        resumen_lines.append(cubre_no_str)
    if bombas_raras_str:
        resumen_lines.append(bombas_raras_str)
    if len(resumen_lines) == 1:
        resumen_lines.append("  · Sin alertas — semana completa y dentro de estándares. ✅")
    resumen_texto = "\n".join(resumen_lines)

    # ── Aplicar sobre plantilla ──
    from pptx.dml.color import RGBColor
    RED_ALERT = RGBColor(0xC0, 0x39, 0x2B)
    ROJOS_S2 = {"Text 69"}   # excepciones lavado → rojo

    # Colores según % (misma paleta que el modelo)
    COL_VERDE  = RGBColor(0x2E, 0x8B, 0x57)  # #2E8B57
    COL_NARAN  = RGBColor(0xE6, 0x7E, 0x22)  # #E67E22
    COL_ROJO   = RGBColor(0xE3, 0x1C, 0x23)  # #E31C23

    def _color_por_pct(pct: float) -> RGBColor:
        if pct >= 80: return COL_VERDE
        if pct >= 50: return COL_NARAN
        return COL_ROJO

    # Cards Slide 1: (nombre_texto_pct, shape_bar_bg, shape_bar_fill, n_ok, total)
    cards_slide1 = [
        ("Text 7",  "Shape 8",  "Shape 9",  n_bomba,  total),
        ("Text 14", "Shape 15", "Shape 16", n_prod,   total),
        ("Text 21", "Shape 22", "Shape 23", n_cera,   total),
        ("Text 28", "Shape 29", "Shape 30", n_shp,    total),
        ("Text 35", "Shape 36", "Shape 37", n_cep,    total),
        ("Text 42", "Shape 43", "Shape 44", n_lavseg, total),
        ("Text 49", "Shape 50", "Shape 51", n_cf_pob, total),
    ]

    def _pintar_card(slide, tpct_name, sbg_name, sfill_name, n_ok, total):
        pct = 100.0 * n_ok / total if total else 0.0
        color = _color_por_pct(pct)
        by_name = {name: sh for name, sh in _walk_named(slide.shapes)}
        # Fondo de la barra: define el ancho total disponible
        bg = by_name.get(sbg_name)
        fill = by_name.get(sfill_name)
        if bg is not None and fill is not None:
            try:
                new_w = max(1, int(bg.width * pct / 100.0))
                fill.left = bg.left    # asegurar alineación izquierda
                fill.top  = bg.top
                fill.width = new_w
                fill.height = bg.height
                fill.fill.solid()
                fill.fill.fore_color.rgb = color
            except Exception as e:
                log(f"  ⚠ pintando {sfill_name}: {e}")
        # Texto del porcentaje: color de la barra
        tpct = by_name.get(tpct_name)
        if tpct is not None and tpct.has_text_frame:
            for p in tpct.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = color

    prs = Presentation(str(_PPT_TEMPLATE))
    slides = list(prs.slides)

    # 1) Reemplazar todos los textos (S1 + S2)
    for idx, (slide, mapping) in enumerate(zip(slides, (replacements_s1, replacements_s2))):
        rojos = ROJOS_S2 if idx == 1 else set()
        for name, sh in _walk_named(slide.shapes):
            if name in mapping:
                col = RED_ALERT if name in rojos else None
                _replace_text(sh, mapping[name], color=col)

    # 2) Pintar barras + color del % en Slide 1
    for args in cards_slide1:
        _pintar_card(slides[0], *args)

    # 3) Ocultar círculos huérfanos en Slide 2 (filas de bomba sin dato)
    if ocultar_shapes_s2:
        for name, sh in _walk_named(slides[1].shapes):
            if name in ocultar_shapes_s2:
                try:
                    # Mover fuera del área visible (más seguro que borrar el XML)
                    sh.left = -10_000_000
                except Exception:
                    pass

    # 4) Ocultar filas vacías del card de PRODUCCIÓN (las que no tienen
    #    datos). El card mismo NO se toca — solo se limpian los textos.
    from pptx.util import Emu, Pt
    n_filas = max(1, filas_prod_activas)
    for name, sh in _walk_named(slides[1].shapes):
        for i in range(n_filas, 6):
            t_n, t_v, t_p = prod_positions[i]
            if name in (t_n, t_v, t_p):
                try:
                    sh.left = -10_000_000
                except Exception:
                    pass

    # 5) Solo mover HACIA ARRIBA los shapes de los cards INFERIORES.
    #    Header, logos y cards superiores intactos.
    DELTA_UP = 500_000     # ~0.55 in (moderado — evita amontonamiento del LAVADO)
    SLIDE_H  = 5_143_500
    # Ajustes finos de textos que se solapaban visualmente:
    #  - Text 63 (Distribución Cera): tap. por "Más frecuente"; lo bajamos
    #  - Text 69 (Excepciones lavado contenido): lo bajamos para separarlo
    #    del subtítulo Text 68.
    _EXTRA_DOWN = {"Text 63": 180_000, "Text 69": 200_000}
    for name, sh in _walk_named(slides[1].shapes):
        try:
            if sh.top is None or sh.top < 0:
                continue
            # Solo aplicar delta a cards INFERIORES (top original >= 3M)
            if sh.top >= 3_000_000:
                new_top = sh.top - DELTA_UP
                if new_top > 0:
                    sh.top = new_top
            # Después de mover, aplicar extra_down si aplica
            if name in _EXTRA_DOWN:
                sh.top = (sh.top or 0) + _EXTRA_DOWN[name]
        except Exception:
            pass

    # 6) Insertar textbox de Resumen en la zona libre abajo (todo el ancho)
    try:
        # Nueva base: bottom de cards inferiores = 3017520 + 1920240 - DELTA_UP
        # (además debe respetar el Text 69 extra_down)
        top_resumen = 3_017_520 + 1_920_240 - DELTA_UP + 220_000
        available_h = SLIDE_H - top_resumen - 60_000
        left_resumen  = 228_600
        width_resumen = 9_144_000 - 2 * 228_600
        box = slides[1].shapes.add_textbox(
            Emu(left_resumen), Emu(top_resumen),
            Emu(width_resumen), Emu(max(available_h, 500_000)))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(50_000)
        tf.margin_top = Emu(20_000); tf.margin_bottom = Emu(0)
        lines_all = resumen_texto.split("\n")
        titulo = lines_all[0]
        items = lines_all[1:]
        # Acotar sub-items de "Faltantes de datos" a los 3 primeros
        limpio = []
        sub_shown = 0
        for ln in items:
            if not ln.strip():
                continue
            if ln.strip().startswith("Faltantes"):
                limpio.append(ln); continue
            if ln.startswith("  ·"):
                if sub_shown < 3:
                    limpio.append(ln); sub_shown += 1
                continue
            limpio.append(ln)
        p0 = tf.paragraphs[0]
        r0 = p0.add_run(); r0.text = titulo
        r0.font.bold = True; r0.font.size = Pt(11)
        r0.font.color.rgb = RGBColor(0x1F, 0x4E, 0x78)
        for ln in limpio:
            p = tf.add_paragraph()
            r = p.add_run(); r.text = ln
            r.font.size = Pt(9)
            if "atípico" in ln.lower() or "≠" in ln or "sin cubre" in ln.lower():
                r.font.color.rgb = RED_ALERT
            elif ln.strip().startswith("Faltantes"):
                r.font.color.rgb = RGBColor(0x1F, 0x4E, 0x78)
                r.font.bold = True
            else:
                r.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
    except Exception as e:
        log(f"  ⚠ No se pudo agregar textbox de resumen: {e}")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), resumen_lines


# ── Email ───────────────────────────────────────────────────────────────────
def _cuerpo_html(mes_lbl: str, sem_iso: int, semana_ini: date, semana_fin: date,
                 n_sem: int, n_mes: int, resumen_lines: list[str]) -> str:
    """Cuerpo HTML del correo. `resumen_lines` = mismas alertas del PPT."""
    # Renderizar el resumen como bloque HTML con las líneas relevantes
    resumen_html = ""
    if resumen_lines and len(resumen_lines) > 1:
        # La primera línea es el título "Resumen — Semana X ..."; el resto son items
        items_html = ""
        for ln in resumen_lines[1:]:
            texto = ln.lstrip(" ·").strip()
            if not texto:
                continue
            # Alertas en rojo
            es_alerta = any(k in texto.lower()
                            for k in ("atípico", "≠", "sin cubre", "sin dato"))
            color = "#c0392b" if es_alerta else "#334155"
            # Sub-items indentados
            if ln.startswith("  ·") or ln.startswith(" ·"):
                items_html += (f'<li style="color:{color};margin:2px 0;">'
                               f'{texto}</li>')
            else:
                items_html += (f'<div style="font-weight:600;color:#334155;'
                               f'margin:6px 0 2px 0;">{texto}</div>')
        if items_html:
            resumen_html = (
                '<div style="background:#f8fafc;border-left:4px solid #1F4E78;'
                'padding:10px 14px;margin:14px 0;border-radius:4px;">'
                '<div style="font-weight:700;color:#1F4E78;margin-bottom:6px;">'
                '📊 Consideraciones importantes de la semana</div>'
                f'<ul style="margin:4px 0 4px 20px;padding:0;font-size:13px;">'
                f'{items_html}</ul>'
                '</div>'
            )

    return f"""\
<!DOCTYPE html>
<html><body style="font-family:Arial,sans-serif;color:#1f2937;max-width:720px;line-height:1.55;">
<p>Estimados, buenos días.</p>

<p>A continuación comparto el resumen consolidado de <b>{mes_lbl}</b>
(hasta el {semana_fin.strftime('%d/%m/%Y')}) de los registros de equipos
y consumos <b>Shell</b>. Semana ISO {sem_iso}.</p>

<ul>
  <li>📎 <b>Presentación (PPT)</b> — consolidado del mes:
      <b>{n_mes}</b> registros analizados.</li>
  <li>📎 <b>Excel</b> con el detalle OT-por-OT del mes
      ({mes_lbl}): <b>{n_mes}</b> OTs.</li>
</ul>

{resumen_html}

<p>Cualquier observación o consulta estaremos atentos.</p>

<p>Saludos,</p>
<p style="margin-bottom:0"><b>Operaciones Occimiano</b><br>
<a href="https://ops-occimiano-dashboard.streamlit.app/" style="color:#2563eb;">
Dashboard de Operaciones</a></p>
</body></html>"""


def _enviar(destinatarios: list, asunto: str, cuerpo_html: str,
            adjuntos: list[tuple[bytes, str, str]],
            in_reply_to: Optional[str], dry_run: bool) -> Optional[str]:
    """Devuelve el Message-ID del correo enviado (o None si dry_run)."""
    if dry_run:
        log(f"  [DRY-RUN] TO: {destinatarios}")
        log(f"  [DRY-RUN] asunto: {asunto}")
        log(f"  [DRY-RUN] in_reply_to: {in_reply_to}")
        for data, fname, ctype in adjuntos:
            log(f"  [DRY-RUN] adjunto: {fname}  ({len(data)//1024} KB)  {ctype}")
        return None

    if not GMAIL_PASSWORD:
        raise RuntimeError("Falta GMAIL_APP_PASSWORD")

    msg = EmailMessage()
    msg["Subject"] = asunto
    msg["From"]    = f"{GMAIL_DISPLAY} <{GMAIL_SENDER}>"
    msg["To"]      = ", ".join(destinatarios)
    msg["Date"]    = formatdate(localtime=True)
    mid = make_msgid(domain="occimiano.cl")
    msg["Message-ID"] = mid
    if in_reply_to:
        msg["In-Reply-To"] = in_reply_to
        msg["References"]  = in_reply_to
    msg.set_content("Este correo contiene HTML. Actualiza tu cliente para verlo.")
    msg.add_alternative(cuerpo_html, subtype="html")
    for data, fname, ctype in adjuntos:
        maintype, subtype = ctype.split("/", 1)
        msg.add_attachment(data, maintype=maintype, subtype=subtype, filename=fname)

    with smtplib.SMTP_SSL("smtp.gmail.com", 465) as s:
        s.login(GMAIL_SENDER, GMAIL_PASSWORD)
        s.send_message(msg)
    log(f"  ✅ Enviado. Message-ID: {mid}")
    return mid


# ── Main ────────────────────────────────────────────────────────────────────
def main() -> int:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    ap = argparse.ArgumentParser()
    ap.add_argument("--dry-run", action="store_true")
    ap.add_argument("--test-email", default=None,
                    help="Redirige a este email en lugar de los destinatarios reales")
    ap.add_argument("--fecha", default=None,
                    help="Simular fecha del lunes de envío (YYYY-MM-DD)")
    ap.add_argument("--force-new-thread", action="store_true",
                    help="Ignora hilo existente y abre uno nuevo")
    ap.add_argument("--mensaje-inicio", default=None,
                    help="Texto HTML que se antepone al cuerpo estándar "
                         "(ej. para 'fe de erratas' o notas puntuales)")
    args = ap.parse_args()

    # Fecha del lunes de envío
    if args.fecha:
        hoy = datetime.strptime(args.fecha, "%Y-%m-%d").date()
    else:
        hoy = (datetime.now(timezone.utc) - timedelta(hours=4)).date()

    # Semana ISO ANTERIOR: cerrada al domingo anterior
    dias_desde_lunes = hoy.weekday()   # 0=lun
    lunes_actual = hoy - timedelta(days=dias_desde_lunes)
    semana_fin = lunes_actual - timedelta(days=1)   # domingo pasado
    semana_ini = semana_fin - timedelta(days=6)     # lunes pasado
    sem_iso = semana_ini.isocalendar().week

    # Mes del HILO = mes del lunes de envío
    mes_yyyy_mm = hoy.strftime("%Y-%m")
    mes_lbl = f"{MESES_ES[hoy.month]} {hoy.year}"

    log(f"═══ REPORTE SHELL CONSUMOS ═══")
    log(f"Lunes envío: {hoy}  ·  Semana anterior: {semana_ini} → {semana_fin} (ISO {sem_iso})")
    log(f"Mes hilo: {mes_yyyy_mm} ({mes_lbl})")

    # ── Gate por día: este correo solo se envía los LUNES ──────────────
    # El workflow corre a diario (los crons semanales de GitHub Actions
    # saltan disparos con frecuencia — best-effort, no garantía).
    # dias_desde_lunes 0 = lunes.
    if (not args.dry_run and not args.test_email and not args.fecha
            and dias_desde_lunes != 0):
        log(f"═══ SKIP: hoy es {hoy.strftime('%A')}, "
            "el correo se envía solo los lunes ═══")
        return 0

    # ── Dedupe idempotente (cron diario + múltiples horas) ──
    from email_dedupe import ya_enviado_hoy, marcar_enviado_hoy
    if not args.dry_run and not args.test_email and ya_enviado_hoy(
            "shell_consumos", fecha=hoy.strftime("%Y-%m-%d")):
        log("═══ SKIP: ya se envió el reporte Shell de hoy (dedupe) ═══")
        return 0

    # 1. Cargar datos (corte = ayer inclusive; hoy se excluye porque puede
    # tener OTs sin cerrar o data aún no sincronizada desde Fracttal)
    corte_hasta = hoy - timedelta(days=1)
    log(f"Corte de datos: hasta {corte_hasta} inclusive (hoy excluido)")
    log("Cargando datos Shell desde Supabase…")
    df_mes, df_sem = load_shell_data(mes_yyyy_mm, semana_ini, semana_fin, corte_hasta)
    log(f"  df_mes: {len(df_mes)} OTs  ·  df_sem: {len(df_sem)} OTs")
    if df_mes.empty:
        log("⚠️  Sin datos Shell en el mes. Salgo.")
        return 0

    # 2. Generar adjuntos
    log("Generando Excel…")
    xlsx = build_excel(df_mes)
    log(f"  Excel: {len(xlsx)//1024} KB")
    log("Generando PPT…")
    # El PPT ahora usa el mes COMPLETO (igual que el Excel). Pasamos
    # semana_ini/semana_fin solo como contexto para el subtítulo pero
    # el análisis se hace sobre df_mes.
    pptx, resumen_lines = build_ppt(df_mes, mes_lbl, sem_iso, semana_ini, semana_fin,
                                     corte_hasta=corte_hasta)
    log(f"  PPT: {len(pptx)//1024} KB · resumen: {len(resumen_lines)} líneas")

    # Nombres de archivo
    fecha_str = hoy.strftime("%d-%m-%Y")
    fname_xlsx = f"Registros Shell {fecha_str}.xlsx"
    fname_pptx = f"Registro_Equipos_Consumo_Shell_{MESES_ES[hoy.month]}{hoy.year}_Sem{sem_iso}.pptx"

    # Dry-run: guardar copia local para revisión
    if args.dry_run:
        out_dir = Path(__file__).parent / "_dry_run_shell"
        out_dir.mkdir(exist_ok=True)
        (out_dir / fname_xlsx).write_bytes(xlsx)
        (out_dir / fname_pptx).write_bytes(pptx)
        log(f"  💾 Guardado local: {out_dir}")

    # 3. Buscar hilo existente del mes
    thread = None if args.force_new_thread else _get_thread(TOPIC, mes_yyyy_mm)
    if thread:
        log(f"  Hilo del mes existe: {thread['message_id']}")
    else:
        log(f"  Nuevo hilo del mes {mes_yyyy_mm}")

    # 4. Asunto
    if thread:
        # Reply: usar el mismo subject con "Re:"
        subject_base = thread.get("subject") or f"Resumen consumos de equipos Shell — {mes_lbl}"
        asunto = subject_base if subject_base.lower().startswith("re:") else f"Re: {subject_base}"
        # Actualizar el cuerpo con la semana actual
        subject_base_for_body = subject_base
    else:
        asunto = f"Resumen consumos de equipos Shell (Sem. {sem_iso}) — {mes_lbl}"
        subject_base_for_body = asunto

    # 5. Destinatarios (test override)
    to_list = [args.test_email] if args.test_email else DESTINATARIOS
    cuerpo = _cuerpo_html(mes_lbl, sem_iso, semana_ini, semana_fin,
                          len(df_sem), len(df_mes), resumen_lines)
    # Prepend "mensaje inicio" (ej. fe de erratas)
    if args.mensaje_inicio:
        cuerpo = args.mensaje_inicio + cuerpo
    if args.test_email:
        asunto = f"[PRUEBA] {asunto}"
        cuerpo = (f'<div style="background:#fef3c7;border-left:4px solid #f59e0b;'
                  f'padding:10px 14px;margin-bottom:14px;color:#78350f;'
                  f'font-family:Arial,sans-serif;">'
                  f'<b>⚠️ MODO PRUEBA</b><br>'
                  f'En envío real este correo iría a: <b>{", ".join(DESTINATARIOS)}</b>'
                  f'</div>' + cuerpo)

    adjuntos = [
        (xlsx, fname_xlsx,
         "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        (pptx, fname_pptx,
         "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    ]

    # 6. Enviar
    try:
        mid = _enviar(to_list, asunto, cuerpo, adjuntos,
                      in_reply_to=(thread["message_id"] if thread else None),
                      dry_run=args.dry_run)
    except Exception as e:
        log(f"❌ ERROR envío: {e}")
        log(traceback.format_exc()[:600])
        return 1

    # 7. Guardar thread si es nuevo hilo (y no test/dry-run)
    if mid and not thread and not args.test_email:
        ok = _save_thread(TOPIC, mes_yyyy_mm, mid, subject_base_for_body)
        log(f"  Thread guardado: {ok}")

    # 8. Marcar envío del día (dedupe backup)
    if mid and not args.test_email:
        marcar_enviado_hoy("shell_consumos",
                           subject_base_for_body,
                           fecha=hoy.strftime("%Y-%m-%d"))

    log("═══ DONE ═══")
    return 0


if __name__ == "__main__":
    sys.exit(main())
